from __future__ import annotations

import io
from pathlib import Path
from typing import Any


_IMAGE_SUFFIXES = {".bmp", ".gif", ".jpeg", ".jpg", ".png", ".tif", ".tiff", ".webp"}
_PLAIN_TEXT_SUFFIXES = {".csv", ".json", ".log", ".md", ".rtf", ".text", ".txt", ".xml", ".yaml", ".yml"}


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding).strip()
        except UnicodeDecodeError:
            continue
    return ""


def _extract_docx_text(data: bytes) -> str:
    try:
        from docx import Document

        document = Document(io.BytesIO(data))
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        table_rows = [
            " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            for table in document.tables
            for row in table.rows
        ]
        return "\n".join(paragraphs + [row for row in table_rows if row]).strip()
    except Exception:
        return ""


def _extract_xlsx_text(data: bytes) -> str:
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        lines: list[str] = []
        for worksheet in workbook.worksheets:
            lines.append(worksheet.title)
            for row in worksheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value not in (None, "")]
                if values:
                    lines.append(" | ".join(values))
        return "\n".join(lines).strip()
    except Exception:
        return ""


def _extract_pptx_text(data: bytes) -> str:
    """Extract readable text from PowerPoint slide shapes and tables."""
    try:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(data))
        lines: list[str] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = "\n".join(
                        paragraph.text.strip()
                        for paragraph in shape.text_frame.paragraphs
                        if paragraph.text.strip()
                    ).strip()
                    if text:
                        slide_lines.append(text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if values:
                            slide_lines.append(" | ".join(values))
            if slide_lines:
                lines.append(f"Slide {slide_number}")
                lines.extend(slide_lines)
        return "\n".join(lines).strip()
    except Exception:
        return ""


def _ocr_image(image: Any) -> str:
    import pytesseract

    return str(pytesseract.image_to_string(image) or "").strip()


def _native_pdf_text_is_sufficient(text: str) -> bool:
    normalized = str(text or "").strip()
    if len(normalized) < 8:
        return False
    alphanumeric_count = sum(character.isalnum() for character in normalized)
    # ponytail: this deliberately cheap heuristic handles normal PDFs; layout-aware detection is the upgrade path.
    return alphanumeric_count / max(len(normalized), 1) >= 0.35


def _ocr_image_with_metadata(image: Any, *, page_number: int = 1) -> dict[str, Any]:
    import pytesseract

    warnings: list[str] = []
    width, height = image.size
    if min(width, height) < 800:
        warnings.append(f"Page {page_number} is low resolution ({width}x{height}); OCR accuracy may be reduced.")

    rotation_degrees = 0
    try:
        osd = pytesseract.image_to_osd(image, output_type=pytesseract.Output.DICT)
        rotation_degrees = int(osd.get("rotate") or 0) % 360
    except Exception:
        rotation_degrees = 0
    if rotation_degrees:
        image = image.rotate(-rotation_degrees, expand=True)

    text = _ocr_image(image)
    confidence_values: list[float] = []
    try:
        ocr_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
        for value in ocr_data.get("conf") or []:
            try:
                confidence = float(value)
            except (TypeError, ValueError):
                continue
            if confidence >= 0:
                confidence_values.append(confidence)
    except Exception:
        confidence_values = []
    confidence = (
        round(sum(confidence_values) / len(confidence_values) / 100, 3)
        if confidence_values
        else (0.5 if text else 0.0)
    )
    if text and confidence < 0.45:
        warnings.append(f"Page {page_number} OCR confidence is low ({confidence:.0%}).")
    return {
        "text": text,
        "page_number": page_number,
        "method": "ocr",
        "status": "ready" if text else "failed",
        "char_count": len(text),
        "confidence": confidence,
        "rotation_degrees": rotation_degrees,
        "width": width,
        "height": height,
        "warnings": warnings,
    }


def _extract_image_text(data: bytes) -> dict[str, Any]:
    try:
        from PIL import Image

        with Image.open(io.BytesIO(data)) as image:
            return _ocr_image_with_metadata(image, page_number=1)
    except Exception:
        return {
            "text": "",
            "page_number": 1,
            "method": "ocr",
            "status": "failed",
            "char_count": 0,
            "confidence": 0.0,
            "rotation_degrees": 0,
            "warnings": ["Image OCR failed before text could be extracted."],
        }


def _extract_pdf_text(data: bytes, *, allow_ocr: bool = True) -> tuple[str, str, list[dict[str, Any]], list[str]]:
    native_pages: list[str] = []
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        native_pages = [str(page.extract_text() or "").strip() for page in reader.pages]
    except Exception:
        native_pages = []

    if native_pages and all(_native_pdf_text_is_sufficient(text) for text in native_pages):
        pages = [
            {
                "page_number": index + 1,
                "method": "native",
                "status": "ready",
                "char_count": len(text),
                "confidence": 1.0,
                "rotation_degrees": 0,
                "warnings": [],
            }
            for index, text in enumerate(native_pages)
        ]
        return "\n".join(native_pages).strip(), "pdf_native", pages, []

    if not allow_ocr:
        pages = [
            {
                "page_number": index + 1,
                "method": "native" if text else "ocr_skipped",
                "status": "ready" if text else "failed",
                "char_count": len(text),
                "confidence": 1.0 if text else 0.0,
                "rotation_degrees": 0,
                "warnings": [] if text else [f"Page {index + 1} needs OCR, but OCR was skipped."],
            }
            for index, text in enumerate(native_pages)
        ]
        return "\n".join(text for text in native_pages if text).strip(), "pdf_ocr_skipped", pages, [
            warning
            for page in pages
            for warning in page.get("warnings") or []
        ]

    try:
        import fitz
        from PIL import Image

        document = fitz.open(stream=data, filetype="pdf")
        pages: list[dict[str, Any]] = []
        page_texts: list[str] = []
        for page_index, page in enumerate(document):
            native_text = native_pages[page_index] if page_index < len(native_pages) else ""
            if _native_pdf_text_is_sufficient(native_text):
                pages.append({
                    "page_number": page_index + 1,
                    "method": "native",
                    "status": "ready",
                    "char_count": len(native_text),
                    "confidence": 1.0,
                    "rotation_degrees": 0,
                    "warnings": [],
                })
                page_texts.append(native_text)
                continue
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
            page_result = _ocr_image_with_metadata(image, page_number=page_index + 1)
            pages.append(page_result)
            page_texts.append(str(page_result.get("text") or ""))
        methods = {str(page.get("method") or "") for page in pages}
        method = "pdf_mixed" if len(methods) > 1 else "pdf_ocr"
        warnings = [warning for page in pages for warning in page.get("warnings") or []]
        return "\n".join(text for text in page_texts if text).strip(), method, pages, warnings
    except Exception:
        return "", "none", [], ["PDF parsing and OCR both failed."]


def extract_document_text(filename: str, data: bytes, *, allow_ocr: bool = True) -> dict[str, Any]:
    suffix = Path(filename or "").suffix.lower()
    warnings: list[str] = []
    pages: list[dict[str, Any]] = []
    text = ""
    method = "none"

    if suffix == ".docx":
        text, method = _extract_docx_text(data), "docx"
    elif suffix == ".pdf":
        text, method, pages, pdf_warnings = _extract_pdf_text(data, allow_ocr=allow_ocr)
        warnings.extend(pdf_warnings)
    elif suffix in _IMAGE_SUFFIXES:
        if allow_ocr:
            image_result = _extract_image_text(data)
            text, method = str(image_result.get("text") or ""), "image_ocr"
            pages = [image_result]
            warnings.extend(image_result.get("warnings") or [])
        else:
            text, method = "", "image_ocr_skipped"
    elif suffix == ".xlsx":
        text, method = _extract_xlsx_text(data), "xlsx"
    elif suffix == ".pptx":
        text, method = _extract_pptx_text(data), "pptx"
    elif suffix in _PLAIN_TEXT_SUFFIXES or not suffix:
        text, method = _decode_text(data), "plain_text"
    else:
        decoded = _decode_text(data)
        if decoded and "\x00" not in decoded:
            text, method = decoded, "plain_text_fallback"

    if not text:
        if suffix == ".pdf" and not allow_ocr:
            warnings.insert(
                0,
                "No embedded PDF text was found and OCR was skipped. "
                "Scanned PDFs need OCR; upload a text-based PDF or DOCX."
            )
        elif suffix in _IMAGE_SUFFIXES and not allow_ocr:
            warnings.append("Image OCR is skipped during upload. Upload a text-based PDF or DOCX.")
        elif suffix in _IMAGE_SUFFIXES or suffix == ".pdf":
            warnings.append("OCR produced no text. Verify that Pillow, pytesseract, PyMuPDF, and Tesseract OCR are installed.")
        else:
            if suffix == ".pptx":
                warnings.append(
                    "PowerPoint text extraction failed. Verify that python-pptx is installed "
                    "and the presentation is a readable .pptx file."
                )
            else:
                warnings.append(f"No text extractor is available for '{suffix or 'unknown'}' files.")
    page_confidences = [
        float(page.get("confidence") or 0)
        for page in pages
        if page.get("status") == "ready"
    ]
    confidence = (
        round(sum(page_confidences) / len(page_confidences), 3)
        if page_confidences
        else (1.0 if text else 0.0)
    )
    failed_pages = [page for page in pages if page.get("status") == "failed"]
    status = "partial" if text and failed_pages else ("ready" if text else "failed")
    return {
        "text": text.strip(),
        "char_count": len(text.strip()),
        "method": method,
        "provider": "",
        "model": "",
        "warnings": list(dict.fromkeys(warnings)),
        "pages": pages,
        "confidence": confidence,
        "status": status,
        "is_ocr": False,
        "is_low_confidence_ocr": False,
        "extracted_at": "",
        "layout_sections": [],
        "experience_details": [],
    }


def create_word_companion_bytes(text: str, *, title: str = "") -> bytes:
    try:
        from docx import Document
    except Exception as exc:
        raise RuntimeError("python-docx is required to create Word CV companions.") from exc

    document = Document()
    if title:
        document.core_properties.title = title
    for line in str(text or "").splitlines():
        document.add_paragraph(line)
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def extraction_metadata(extraction: dict[str, Any]) -> dict[str, Any]:
    return {
        "source_text": str(extraction.get("text") or ""),
        "source_char_count": int(extraction.get("char_count") or 0),
        "text_extraction": {
            "method": str(extraction.get("method") or "none"),
            "provider": str(extraction.get("provider") or ""),
            "model": str(extraction.get("model") or ""),
            "warnings": list(extraction.get("warnings") or []),
            "status": str(extraction.get("status") or "failed"),
            "confidence": float(extraction.get("confidence") or 0),
            "pages": list(extraction.get("pages") or []),
            "is_ocr": bool(extraction.get("is_ocr")),
            "is_low_confidence_ocr": bool(extraction.get("is_low_confidence_ocr")),
            "extracted_at": str(extraction.get("extracted_at") or ""),
            "layout_sections": list(extraction.get("layout_sections") or []),
            "experience_details": list(extraction.get("experience_details") or []),
        },
    }


__all__ = ["create_word_companion_bytes", "extract_document_text", "extraction_metadata"]
